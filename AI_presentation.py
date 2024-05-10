import requests
import json
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup
import textwrap
from vscode_openai_config import OPENAI_API_KEY

def generate_text_with_gpt(prompt):
    api_key = OPENAI_API_KEY
    headers = {"Authorization": f"Bearer {api_key}"}
    data = {
        "model": "text-davinci-003",
        "prompt": prompt,
        "temperature": 0.7,
        "max_tokens": 250,
        "top_p": 1.0,
        "frequency_penalty": 0.0,
        "presence_penalty": 0.0
    }
    response = requests.post("https://api.openai.com/v1/completions", headers=headers, json=data)
    if response.status_code == 200:
        response_json = response.json()
        return response_json["choices"][0]["text"].strip()
    else:
        print(f"Failed to fetch data from OpenAI. Status code: {response.status_code}")
        return ""

def find_image_urls(query):
    url = f"https://unsplash.com/s/photos/{query.replace(' ', '-')}"
    response = requests.get(url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.text, 'html.parser')
        images = soup.find_all('img', {'class': '_2zEKz'})
        urls = [img['src'] for img in images[:1]]  # Fetch only the first image for simplicity
        return urls
    return []

def add_slide_with_text_and_image(presentation, title, content, image_url):
    slide_layout = presentation.slide_layouts[5]  # Use a blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set title manually if needed
    title_shape = slide.shapes.title
    title_shape.text = title

    # Add a textbox for content instead of using placeholders
    inches_from_top = 2
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(inches_from_top), Inches(9), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.clear()

    # Splitting content to manage overflow
    for line in split_text_into_lines(content, max_width=75):  # Adjust max_width as needed
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)  # Adjust font size as needed

    # Add image
    if image_url:
        response = requests.get(image_url)
        image_stream = io.BytesIO(response.content)
        # Adjust the position and size of the image as needed
        slide.shapes.add_picture(image_stream, Inches(0.5), Inches(inches_from_top + 3), width=Inches(6))

def split_text_into_lines(text, max_width=75):
    """
    Splits text into lines, trying not to exceed the specified max_width of characters.
    This is a simple approach and might need adjustments based on your specific font and layout.
    """
    lines = textwrap.wrap(text, width=max_width)
    return lines

def main():
    topic = "humonoid robots"  # Adjust the topic as needed
    presentation = Presentation()

    aspects = ["introduction", "history", "applications", "future"]
    for aspect in aspects:
        prompt = f"Explain the {aspect} of {topic} in detail."
        content = generate_text_with_gpt(prompt)
        image_urls = find_image_urls(f"{topic} {aspect}")
        image_url = image_urls[0] if image_urls else None
        add_slide_with_text_and_image(presentation, aspect.capitalize(), content, image_url)

    presentation.save(f"{topic.replace(' ', '_')}_presentation.pptx")
    print(f"Presentation for {topic} created successfully.")

if __name__ == "__main__":
    main()
